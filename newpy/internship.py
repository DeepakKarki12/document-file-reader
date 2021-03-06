from gtts import gTTS
import streamlit as st
import pdfplumber
import docx2txt
import pandas as pd
import time
import pip
from pptx import Presentation



pip.main(["install", "openpyxl"])
counter=0


def convert_to_audio(text):
    global counter
    counter +=1
    audio = gTTS(text=text, lang='en', tld='com')
    audio.save(f"textaud{counter}.mp3")
    audio_file = open(f"textaud{counter}.mp3", "rb")
    st.audio(audio_file.read())
    with open(f"textaud{counter}.mp3", 'rb') as f:
        st.download_button('Download Mp3', f, file_name=f"textaud{counter}.mp3")


count = 0
interface = ""
total_text = ""
raw_text = ""
# st.title("Document Reader")
menu = ["DocumentFiles","ExcelFile","PPT"]
choice = st.sidebar.selectbox("Menu",menu)

if choice == "DocumentFiles":
    st.title("DocumentFiles")
    docx_file = st.file_uploader("Upload Document", type=["pdf","docx","txt"])
    if docx_file is not None:
        interface = "File has been uploaded"
        st.subheader(interface)
        if docx_file.type == "application/pdf":
            starting_page_no = st.number_input('enter the starting page number of document',value= 1)
            ending_page_no = st.number_input('enter the ending page number of document', value=1)
            total_pages = ending_page_no - starting_page_no
            if st.button("Process"):
                for i in range(total_pages + 1):
                    loop = 1
                    try:
                        with pdfplumber.open(docx_file) as pdf:
                            pages = pdf.pages[int(starting_page_no - 1 + i)]
                            line = pages.extract_text()
                            total_text = total_text + line
                            if len(total_text) < 40000:
                                pass
                            else:
                                convert_to_audio(total_text)
                                time.sleep(999)
                                total_text = ""
                                loop = 0

                    except:
                        st.write("None")
                    time.sleep(2)
                if loop == 1:
                    convert_to_audio(total_text)
#                 st.write("If you last converted audio file were more than 1 hour than please wait for 10 minute")
#                 st.write(len(total_text))

        elif docx_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            if st.button("Process"):
#                 file_details = {"filename": docx_file.name, "filetype": docx_file.type,
#                                 "filesize": docx_file.size}
#                 st.write(file_details)
                total_text = docx2txt.process(docx_file)
                convert_to_audio(total_text)
                # audio_file = open("textaud.mp3", "rb")
                # st.audio(audio_file.read())
        elif docx_file.type == "text/plain":
            if st.button("Process"):
#                 file_details = {"filename": docx_file.name, "filetype": docx_file.type,
#                                 "filesize": docx_file.size}
#                 st.write(file_details)
                total_text = str(docx_file.read(), "utf-8")
                convert_to_audio(total_text)
    elif docx_file is None:
        st.subheader(interface)

if choice == "ExcelFile":
    st.title("Dataset")
    data_file = st.file_uploader("Upload CSV", type=["xlsx"])
    if data_file is not None:
        interface = "file has been uploaded"
        st.subheader(interface)
        if st.button("Process"):
#             file_details = {"filename": data_file.name, "filetype": data_file.type,
#                             "filesize": data_file.size}
            xl = pd.ExcelFile(data_file)
            for sheet in xl.sheet_names:
                file = pd.read_excel(xl, sheet_name=sheet)
                docx_file = file.to_csv(sheet + '.txt', header=False, index=False)
                docx_file = open('Sheet1.txt','rb')
                total_text = str(docx_file.read(), "utf-8")
                convert_to_audio(total_text)

    else:
        st.subheader(interface)

if choice == "PPT":
    st.title("PPT")
    data_file = st.file_uploader("Upload CSV", type=["pptx"])
    if data_file is not None:
        interface = "file has been uploaded"
        st.subheader(interface)
#         file_details = {"filename": data_file.name, "filetype": data_file.type,
#                         "filesize": data_file.size}
        final_text = ""
        # for eachfile in glob.glob("*pptx"):
        prs = Presentation(data_file)
        print(data_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    total_text += shape.text
        print(total_text)
        if st.button("Process"):
            convert_to_audio(total_text)
    else:
        st.subheader(interface)
